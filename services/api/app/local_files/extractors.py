from __future__ import annotations

import csv
import hashlib
import io
import json
import logging
import re
import time
import unicodedata
import zipfile
from dataclasses import dataclass, field, replace
from html.parser import HTMLParser
from pathlib import Path
from typing import Final, Iterable, Mapping

from defusedxml import ElementTree as SafeElementTree

logger = logging.getLogger(__name__)


TEXT_EXTENSIONS: Final[frozenset[str]] = frozenset(
    {
        ".txt", ".md", ".rst", ".log", ".py", ".js", ".jsx", ".ts", ".tsx",
        ".java", ".c", ".cpp", ".h", ".hpp", ".cs", ".go", ".rs", ".html",
        ".htm", ".css", ".scss", ".sql", ".sh", ".ps1", ".bat", ".cmd",
        ".toml", ".ini", ".cfg", ".yaml", ".yml", ".jsonl",
    }
)
CONTENT_EXTENSIONS: Final[frozenset[str]] = TEXT_EXTENSIONS | frozenset(
    {".json", ".xml", ".csv", ".pdf", ".docx", ".pptx", ".xlsx"}
)
LEGACY_METADATA_ONLY: Final[frozenset[str]] = frozenset({".doc", ".xls", ".ppt"})


@dataclass(frozen=True)
class ExtractionLimits:
    max_file_size: int = 25 * 1024 * 1024
    max_extracted_chars: int = 2_000_000
    max_unit_chars: int = 8_000
    max_pdf_pages: int = 500
    max_spreadsheet_rows: int = 20_000
    max_spreadsheet_cells: int = 200_000
    spreadsheet_row_group: int = 50
    max_office_uncompressed: int = 100 * 1024 * 1024
    max_office_members: int = 10_000
    # Part 11.1: optional OcrSettings. None means "read from environment", so
    # every existing caller and every Part 11 test keeps its behaviour.
    ocr: "object | None" = None


@dataclass(frozen=True)
class ExtractedUnit:
    text: str
    page_number: int | None = None
    slide_number: int | None = None
    sheet_name: str | None = None
    row_start: int | None = None
    row_end: int | None = None
    line_start: int | None = None
    line_end: int | None = None
    section: str | None = None
    extraction_method: str = "native"
    confidence: float = 1.0
    truncated: bool = False


@dataclass(frozen=True)
class ExtractionResult:
    parser_name: str
    parser_version: str
    status: str
    units: tuple[ExtractedUnit, ...] = ()
    needs_ocr: bool = False
    error_code: str | None = None
    truncated: bool = False
    # Part 11.1 OCR provenance. Defaults keep every Part 11 construction valid.
    ocr_pages: tuple[int, ...] = ()
    pages_needing_ocr: tuple[int, ...] = ()
    ocr_availability: str = "disabled"
    ocr_stats: Mapping[str, object] = field(default_factory=dict)

    @property
    def text(self) -> str:
        return "\n".join(unit.text for unit in self.units)

    @property
    def used_ocr(self) -> bool:
        return bool(self.ocr_pages)


class _HTMLText(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []

    def handle_data(self, data: str) -> None:
        if data.strip():
            self.parts.append(data.strip())


def _sanitize_text(value: object) -> str:
    text = str(value or "").replace("\r\n", "\n").replace("\r", "\n")
    return "".join(
        character
        for character in text
        if character in "\n\t" or unicodedata.category(character) not in {"Cc", "Cs"}
    )


def _decode_text(data: bytes) -> str:
    if data.startswith((b"\xff\xfe", b"\xfe\xff")):
        return data.decode("utf-16")
    if data.startswith(b"\xef\xbb\xbf"):
        return data.decode("utf-8-sig")
    sample = data[:8192]
    if b"\x00" in sample:
        raise ValueError("binary_content")
    control = sum(byte < 9 or (13 < byte < 32) for byte in sample)
    if sample and control / len(sample) > 0.03:
        raise ValueError("binary_content")
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode("cp1252")


def _read_text(path: Path, limits: ExtractionLimits) -> str:
    with path.open("rb") as stream:
        data = stream.read(limits.max_file_size + 1)
    if len(data) > limits.max_file_size:
        raise OverflowError("file_too_large")
    return _sanitize_text(_decode_text(data))


def _line_units(text: str, limits: ExtractionLimits, *, method: str = "native") -> list[ExtractedUnit]:
    lines = text.splitlines()
    units: list[ExtractedUnit] = []
    start = 0
    buffer: list[str] = []
    size = 0
    for index, line in enumerate(lines, 1):
        if not buffer:
            start = index
        if buffer and (size + len(line) + 1 > limits.max_unit_chars or (not line.strip() and size > 500)):
            units.append(ExtractedUnit("\n".join(buffer).strip(), line_start=start, line_end=index - 1, extraction_method=method))
            buffer, size, start = [], 0, index
        buffer.append(line)
        size += len(line) + 1
    if buffer:
        units.append(ExtractedUnit("\n".join(buffer).strip(), line_start=start, line_end=len(lines), extraction_method=method))
    return [unit for unit in units if unit.text]


def _safe_office_package(path: Path, limits: ExtractionLimits) -> None:
    with zipfile.ZipFile(path) as archive:
        members = archive.infolist()
        if len(members) > limits.max_office_members:
            raise OverflowError("office_member_limit")
        if sum(item.file_size for item in members) > limits.max_office_uncompressed:
            raise OverflowError("office_uncompressed_limit")


def _extract_pdf(path: Path, limits: ExtractionLimits) -> ExtractionResult:
    """Extract PDF text page by page, OCRing only pages that genuinely need it.

    Part 11.1. The decision is made PER PAGE, so a mixed PDF keeps its good
    native pages verbatim and OCRs only the image-only ones, with page
    provenance preserved either way.

    Native text always wins. OCR runs only when a page's native text is too
    sparse to trust AND OCR is actually available. When OCR is unavailable,
    disabled, out of budget or fails, the page stays honestly marked needs_ocr:
    no text is ever fabricated, and the PDF itself is never modified.
    """
    import pymupdf

    from . import ocr as ocr_module

    settings = limits.ocr or ocr_module.settings_from_env()
    status = ocr_module.ocr_status(settings)
    budget = ocr_module.OcrBudget(settings)

    units: list[ExtractedUnit] = []
    pages_needing_ocr: list[int] = []

    with pymupdf.open(path) as document:
        page_count = min(len(document), limits.max_pdf_pages)
        for page_index in range(page_count):
            page_number = page_index + 1
            page = document[page_index]
            blocks = page.get_text("blocks", sort=True)
            native_text = _sanitize_text(
                "\n".join(
                    str(block[4]).strip()
                    for block in blocks
                    if len(block) > 4 and str(block[4]).strip()
                )
            )

            if not ocr_module.page_needs_ocr(native_text, settings):
                for unit in _line_units(native_text, limits, method="pymupdf_native"):
                    units.append(replace(unit, page_number=page_number))
                continue

            # This page's native text is insufficient. Try OCR if we can.
            if not status.available or not budget.allows_another_page():
                pages_needing_ocr.append(page_number)
                # Keep whatever sparse native text exists rather than dropping
                # it, but the page is still reported as needing OCR.
                for unit in _line_units(native_text, limits, method="pymupdf_native"):
                    units.append(replace(unit, page_number=page_number))
                continue

            started = time.monotonic()
            try:
                ocr_text = _sanitize_text(
                    ocr_module.ocr_page_text(page, settings, status)
                )
            except ocr_module.OcrError as exc:
                budget.record(succeeded=False, seconds=time.monotonic() - started)
                logger.info("OCR skipped page %s of %s: %s", page_number, path.name, exc)
                pages_needing_ocr.append(page_number)
                for unit in _line_units(native_text, limits, method="pymupdf_native"):
                    units.append(replace(unit, page_number=page_number))
                continue

            budget.record(succeeded=True, seconds=time.monotonic() - started)
            for unit in _line_units(
                ocr_text, limits, method=ocr_module.OCR_EXTRACTION_METHOD
            ):
                units.append(
                    replace(
                        unit,
                        page_number=page_number,
                        confidence=ocr_module.OCR_CONFIDENCE,
                    )
                )

        truncated = len(document) > page_count

    return ExtractionResult(
        "pymupdf",
        pymupdf.__version__,
        "indexed",
        tuple(units),
        # Honest: true when at least one page still lacks trustworthy text.
        needs_ocr=bool(pages_needing_ocr),
        truncated=truncated,
        ocr_pages=tuple(
            page.page_number
            for page in units
            if page.extraction_method == ocr_module.OCR_EXTRACTION_METHOD
            and page.page_number is not None
        ),
        pages_needing_ocr=tuple(pages_needing_ocr),
        ocr_availability=status.availability,
        ocr_stats=budget.as_dict(),
    )


def _extract_docx(path: Path, limits: ExtractionLimits) -> ExtractionResult:
    import docx
    from docx import Document

    _safe_office_package(path, limits)
    document = Document(path)
    units: list[ExtractedUnit] = []
    section: str | None = None
    for paragraph in document.paragraphs:
        text = _sanitize_text(paragraph.text).strip()
        if not text:
            continue
        style = str(getattr(paragraph.style, "name", "") or "")
        if style.casefold().startswith("heading"):
            section = text
        units.append(ExtractedUnit(text[: limits.max_unit_chars], section=section, extraction_method="python-docx", truncated=len(text) > limits.max_unit_chars))
    for table_number, table in enumerate(document.tables, 1):
        rows = [" | ".join(_sanitize_text(cell.text).strip() for cell in row.cells) for row in table.rows]
        text = "\n".join(row for row in rows if row.strip(" |"))
        if text:
            units.append(ExtractedUnit(text[: limits.max_unit_chars], section=section or f"table:{table_number}", extraction_method="python-docx", truncated=len(text) > limits.max_unit_chars))
    return ExtractionResult("python-docx", docx.__version__, "indexed", tuple(units))


def _shape_texts(shape: object) -> Iterable[str]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            yield from _shape_texts(child)
    if getattr(shape, "has_text_frame", False):
        text = _sanitize_text(getattr(shape, "text", "")).strip()
        if text:
            yield text
    if getattr(shape, "has_table", False):
        rows = [" | ".join(_sanitize_text(cell.text).strip() for cell in row.cells) for row in shape.table.rows]
        text = "\n".join(row for row in rows if row.strip(" |"))
        if text:
            yield text


def _extract_pptx(path: Path, limits: ExtractionLimits) -> ExtractionResult:
    import pptx
    from pptx import Presentation

    _safe_office_package(path, limits)
    presentation = Presentation(path)
    units: list[ExtractedUnit] = []
    for slide_number, slide in enumerate(presentation.slides, 1):
        text = "\n".join(part for shape in slide.shapes for part in _shape_texts(shape))
        if text:
            units.append(ExtractedUnit(text[: limits.max_unit_chars], slide_number=slide_number, extraction_method="python-pptx", truncated=len(text) > limits.max_unit_chars))
    return ExtractionResult("python-pptx", pptx.__version__, "indexed", tuple(units))


def _extract_xlsx(path: Path, limits: ExtractionLimits) -> ExtractionResult:
    import openpyxl

    _safe_office_package(path, limits)
    workbook = openpyxl.load_workbook(path, read_only=True, data_only=True, keep_links=False)
    units: list[ExtractedUnit] = []
    rows_seen = cells_seen = 0
    truncated = False
    try:
        for sheet in workbook.worksheets:
            group: list[str] = []
            group_start = 1
            group_end = 0
            for row_number, row in enumerate(sheet.iter_rows(values_only=True), 1):
                if rows_seen >= limits.max_spreadsheet_rows or cells_seen >= limits.max_spreadsheet_cells:
                    truncated = True
                    break
                values = [str(value) for value in row if value is not None]
                rows_seen += 1
                cells_seen += len(row)
                if values:
                    if not group:
                        group_start = row_number
                    group.append(" | ".join(values)[: limits.max_unit_chars])
                    group_end = row_number
                if len(group) >= limits.spreadsheet_row_group:
                    text = "\n".join(group)[: limits.max_unit_chars]
                    units.append(ExtractedUnit(text, sheet_name=sheet.title, row_start=group_start, row_end=group_end, extraction_method="openpyxl", truncated=len("\n".join(group)) > limits.max_unit_chars))
                    group = []
            if group:
                text = "\n".join(group)
                units.append(ExtractedUnit(text[: limits.max_unit_chars], sheet_name=sheet.title, row_start=group_start, row_end=group_end, extraction_method="openpyxl", truncated=len(text) > limits.max_unit_chars))
            if truncated:
                break
    finally:
        workbook.close()
    return ExtractionResult("openpyxl", openpyxl.__version__, "indexed", tuple(units), truncated=truncated)


def _extract_csv(path: Path, limits: ExtractionLimits) -> ExtractionResult:
    text = _read_text(path, limits)
    csv.field_size_limit(min(limits.max_unit_chars, 1_000_000))
    reader = csv.reader(io.StringIO(text))
    units: list[ExtractedUnit] = []
    group: list[str] = []
    start = 1
    row_end = 0
    for row_number, row in enumerate(reader, 1):
        if row_number > limits.max_spreadsheet_rows:
            break
        rendered = " | ".join(field[: limits.max_unit_chars] for field in row)
        if not group:
            start = row_number
        group.append(rendered)
        row_end = row_number
        if len(group) >= limits.spreadsheet_row_group:
            value = "\n".join(group)
            units.append(ExtractedUnit(value[: limits.max_unit_chars], row_start=start, row_end=row_end, extraction_method="csv", truncated=len(value) > limits.max_unit_chars))
            group = []
    if group:
        value = "\n".join(group)
        units.append(ExtractedUnit(value[: limits.max_unit_chars], row_start=start, row_end=row_end, extraction_method="csv", truncated=len(value) > limits.max_unit_chars))
    return ExtractionResult("csv", "stdlib", "indexed", tuple(units), truncated=row_end >= limits.max_spreadsheet_rows)


def _flatten_json(value: object, *, prefix: str = "", output: list[str], limit: int) -> None:
    if sum(len(item) for item in output) >= limit:
        return
    if isinstance(value, dict):
        for key, child in value.items():
            _flatten_json(child, prefix=f"{prefix}.{key}" if prefix else str(key), output=output, limit=limit)
    elif isinstance(value, list):
        for index, child in enumerate(value):
            _flatten_json(child, prefix=f"{prefix}[{index}]", output=output, limit=limit)
    else:
        output.append(f"{prefix}: {value}")


def _bounded_result(result: ExtractionResult, limits: ExtractionLimits) -> ExtractionResult:
    units: list[ExtractedUnit] = []
    remaining = limits.max_extracted_chars
    truncated = result.truncated
    for unit in result.units:
        if remaining <= 0:
            truncated = True
            break
        text = unit.text[:remaining]
        units.append(replace(unit, text=text, truncated=unit.truncated or len(text) < len(unit.text)))
        remaining -= len(text)
    return replace(result, units=tuple(units), truncated=truncated)


def extract_file(path: Path | str, limits: ExtractionLimits | None = None) -> ExtractionResult:
    """Extract one already-policy-approved file without ever modifying it."""
    target = Path(path)
    policy = limits or ExtractionLimits()
    try:
        size = target.stat().st_size
        if size > policy.max_file_size:
            return ExtractionResult("none", "1", "metadata_only", error_code="file_too_large")
        extension = target.suffix.casefold()
        if extension not in CONTENT_EXTENSIONS:
            return ExtractionResult("none", "1", "metadata_only", error_code="unsupported_format")
        if extension == ".pdf":
            result = _extract_pdf(target, policy)
        elif extension == ".docx":
            result = _extract_docx(target, policy)
        elif extension == ".pptx":
            result = _extract_pptx(target, policy)
        elif extension == ".xlsx":
            result = _extract_xlsx(target, policy)
        elif extension == ".csv":
            result = _extract_csv(target, policy)
        else:
            text = _read_text(target, policy)
            if extension == ".json":
                output: list[str] = []
                _flatten_json(json.loads(text), output=output, limit=policy.max_extracted_chars)
                text = "\n".join(output)
            elif extension == ".xml":
                root = SafeElementTree.fromstring(text)
                text = "\n".join(part.strip() for part in root.itertext() if part.strip())
            elif extension in {".html", ".htm"}:
                parser = _HTMLText()
                parser.feed(text)
                text = "\n".join(parser.parts)
            result = ExtractionResult("text", "1", "indexed", tuple(_line_units(text, policy)))
        return _bounded_result(result, policy)
    except OverflowError as exc:
        return ExtractionResult("none", "1", "metadata_only", error_code=str(exc))
    except Exception:
        return ExtractionResult("error", "1", "error", error_code="parse_error")


def unit_hash(unit: ExtractedUnit) -> str:
    return hashlib.sha256(unit.text.encode("utf-8")).hexdigest()


__all__ = [
    "CONTENT_EXTENSIONS",
    "ExtractionLimits",
    "ExtractionResult",
    "ExtractedUnit",
    "LEGACY_METADATA_ONLY",
    "TEXT_EXTENSIONS",
    "extract_file",
    "unit_hash",
]
