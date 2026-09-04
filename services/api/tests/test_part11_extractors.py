from __future__ import annotations

import hashlib
import tempfile
import unittest
from pathlib import Path

import pymupdf
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from services.api.app.local_files.extractors import ExtractionLimits, extract_file
from services.api.app.local_files.redaction import REDACTION_MARKER, redact_secrets


class Part11ExtractorTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temp = tempfile.TemporaryDirectory()
        self.root = Path(self.temp.name)

    def tearDown(self) -> None:
        self.temp.cleanup()

    def test_text_structured_and_multilingual_formats(self) -> None:
        cases = {
            "notes.txt": "English vector database\nहिन्दी सामग्री\nગુજરાતી સામગ્રી",
            "data.json": '{"topic":"retrieval", "count":2}',
            "config.yaml": "topic: retrieval\nitems:\n  - local",
            "safe.xml": "<root><topic>retrieval</topic></root>",
            "rows.csv": "name,topic\nParth,vector database\n",
        }
        expected = {
            "notes.txt": "ગુજરાતી",
            "data.json": "retrieval",
            "config.yaml": "retrieval",
            "safe.xml": "retrieval",
            "rows.csv": "vector database",
        }
        for name, content in cases.items():
            with self.subTest(name=name):
                path = self.root / name
                path.write_text(content, encoding="utf-8")
                result = extract_file(path)
                self.assertEqual(result.status, "indexed")
                self.assertTrue(result.units)
                self.assertIn(expected[name], result.text)

    def test_pdf_page_provenance_and_blank_pdf_ocr_detection(self) -> None:
        native = self.root / "Resume_Parth.pdf"
        doc = pymupdf.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Vector database resume experience")
        doc.save(native)
        doc.close()
        result = extract_file(native)
        self.assertEqual(result.status, "indexed")
        self.assertFalse(result.needs_ocr)
        self.assertEqual(result.units[0].page_number, 1)

        scanned = self.root / "scanned_like.pdf"
        doc = pymupdf.open()
        doc.new_page()
        doc.save(scanned)
        doc.close()
        result = extract_file(scanned)
        self.assertTrue(result.needs_ocr)
        self.assertEqual(result.status, "indexed")

    def test_docx_paragraph_heading_and_table(self) -> None:
        path = self.root / "AI_Architecture.docx"
        doc = Document()
        doc.add_heading("Retrieval", level=1)
        doc.add_paragraph("Bunnelby architecture uses local lexical search.")
        table = doc.add_table(rows=1, cols=2)
        table.cell(0, 0).text = "Engine"
        table.cell(0, 1).text = "FTS5"
        doc.save(path)
        result = extract_file(path)
        self.assertIn("FTS5", result.text)
        self.assertTrue(any(unit.section == "Retrieval" for unit in result.units))

    def test_pptx_slide_and_xlsx_sheet_row_provenance(self) -> None:
        deck = self.root / "presentation.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Local Search"
        slide.placeholders[1].text = "Deterministic retrieval"
        presentation.save(deck)
        deck_result = extract_file(deck)
        self.assertIn("Deterministic retrieval", deck_result.text)
        self.assertTrue(all(unit.slide_number == 1 for unit in deck_result.units))

        book_path = self.root / "budget.xlsx"
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Budget"
        sheet.append(("Item", "Amount"))
        sheet.append(("GPU", 500))
        sheet["C2"] = "=SUM(B2:B2)"
        workbook.save(book_path)
        workbook.close()
        book_result = extract_file(book_path)
        self.assertIn("GPU", book_result.text)
        self.assertNotIn("=SUM", book_result.text)
        self.assertEqual(book_result.units[0].sheet_name, "Budget")
        self.assertEqual(book_result.units[0].row_start, 1)

    def test_corrupt_file_isolated_and_size_cap_is_metadata_only(self) -> None:
        corrupt = self.root / "corrupt.pdf"
        corrupt.write_bytes(b"not a pdf")
        result = extract_file(corrupt)
        self.assertEqual(result.status, "error")
        self.assertEqual(result.error_code, "parse_error")

        large = self.root / "large.log"
        large.write_bytes(b"x" * 200)
        result = extract_file(large, ExtractionLimits(max_file_size=100))
        self.assertEqual(result.status, "metadata_only")
        self.assertEqual(result.error_code, "file_too_large")

    def test_original_bytes_are_unchanged(self) -> None:
        path = self.root / "immutable.txt"
        path.write_text("read only fixture", encoding="utf-8")
        before = hashlib.sha256(path.read_bytes()).digest()
        extract_file(path)
        after = hashlib.sha256(path.read_bytes()).digest()
        self.assertEqual(before, after)

    def test_embedded_credentials_are_redacted_conservatively(self) -> None:
        source = (
            "password = super-secret-value\n"
            "Authorization: Bearer abcdefghijklmnopqrstuvwxyz123456\n"
            "github_token=ghp_abcdefghijklmnopqrstuvwxyz1234567890\n"
            "ordinary prose about password managers remains useful"
        )
        redacted, count = redact_secrets(source)
        self.assertGreaterEqual(count, 3)
        self.assertNotIn("super-secret-value", redacted)
        self.assertNotIn("ghp_", redacted)
        self.assertIn(REDACTION_MARKER, redacted)
        self.assertIn("ordinary prose", redacted)


if __name__ == "__main__":
    unittest.main()
